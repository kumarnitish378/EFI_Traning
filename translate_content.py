#!/usr/bin/env python3
"""
Translate content of files from Spanish to English
"""

import os
from pathlib import Path
from deep_translator import GoogleTranslator
import time

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

TARGET_DIR = Path('/workspaces/EFI_EnglishTraining')
BATCH_SIZE = 5  # Process 5 files before pausing to avoid rate limits

def translate_text(text, max_length=4500):
    """Translate Spanish text to English with chunking for long texts"""
    if not text or not text.strip():
        return text
    
    try:
        translator = GoogleTranslator(source_language='es', target_language='en')
        
        if len(text) > max_length:
            # Split into sentences or words
            chunks = []
            current_chunk = ""
            
            for sentence in text.split('.'):
                if len(current_chunk) + len(sentence) + 1 > max_length:
                    if current_chunk:
                        chunks.append(current_chunk.strip() + '.')
                    current_chunk = sentence
                else:
                    if current_chunk:
                        current_chunk += '. ' + sentence
                    else:
                        current_chunk = sentence
            
            if current_chunk:
                chunks.append(current_chunk.strip())
            
            translated_chunks = []
            for chunk in chunks:
                try:
                    if chunk.strip():
                        translated_chunks.append(translator.translate(chunk))
                except Exception as e:
                    translated_chunks.append(chunk)
            
            return ' '.join(translated_chunks)
        else:
            return translator.translate(text)
    except Exception as e:
        print(f"  Translation error: {str(e)[:100]}")
        return text

def translate_docx(file_path):
    """Translate DOCX file content"""
    if not HAS_DOCX:
        return False
    
    try:
        doc = Document(file_path)
        translated_count = 0
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                original = paragraph.text
                translated = translate_text(original)
                if translated != original:
                    paragraph.text = translated
                    translated_count += 1
        
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            original = paragraph.text
                            translated = translate_text(original)
                            if translated != original:
                                paragraph.text = translated
                                translated_count += 1
        
        if translated_count > 0:
            doc.save(file_path)
            print(f"  ✓ Translated {file_path.name} ({translated_count} items)")
        return True
    except Exception as e:
        print(f"  ✗ Error with {file_path.name}: {str(e)[:100]}")
        return False

def translate_xlsx(file_path):
    """Translate XLSX file content"""
    if not HAS_XLSX:
        return False
    
    try:
        wb = openpyxl.load_workbook(file_path)
        translated_count = 0
        
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str) and cell.value.strip():
                        original = cell.value
                        translated = translate_text(original)
                        if translated != original:
                            cell.value = translated
                            translated_count += 1
        
        if translated_count > 0:
            wb.save(file_path)
            print(f"  ✓ Translated {file_path.name} ({translated_count} items)")
        return True
    except Exception as e:
        print(f"  ✗ Error with {file_path.name}: {str(e)[:100]}")
        return False

def translate_pptx(file_path):
    """Translate PPTX file content"""
    if not HAS_PPTX:
        return False
    
    try:
        prs = Presentation(file_path)
        translated_count = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.text.strip():
                                    translated = translate_text(run.text)
                                    if translated != run.text:
                                        run.text = translated
                                        translated_count += 1
        
        if translated_count > 0:
            prs.save(file_path)
            print(f"  ✓ Translated {file_path.name} ({translated_count} items)")
        return True
    except Exception as e:
        print(f"  ✗ Error with {file_path.name}: {str(e)[:100]}")
        return False

def process_files():
    """Process all files in the target directory"""
    docx_files = list(TARGET_DIR.rglob('*.docx'))
    xlsx_files = list(TARGET_DIR.rglob('*.xlsx'))
    pptx_files = list(TARGET_DIR.rglob('*.pptx'))
    
    all_files = docx_files + xlsx_files + pptx_files
    
    # Filter out files from subdirectories we want to skip
    skip_names = ['copy_and_rename.py', 'translate_project.py', 'main.py']
    all_files = [f for f in all_files if f.name not in skip_names]
    
    print(f"Found {len(docx_files)} .docx files")
    print(f"Found {len(xlsx_files)} .xlsx files")
    print(f"Found {len(pptx_files)} .pptx files\n")
    
    processed = 0
    for idx, file_path in enumerate(all_files):
        print(f"[{idx+1}/{len(all_files)}] Processing {file_path.name}...")
        
        if file_path.suffix.lower() == '.docx':
            translate_docx(file_path)
        elif file_path.suffix.lower() == '.xlsx':
            translate_xlsx(file_path)
        elif file_path.suffix.lower() == '.pptx':
            translate_pptx(file_path)
        
        processed += 1
        
        # Add delay to avoid rate limiting
        if processed % BATCH_SIZE == 0 and idx < len(all_files) - 1:
            print(f"  Waiting 5 seconds to avoid rate limits...\n")
            time.sleep(5)

def main():
    print("Starting content translation...")
    print(f"Target directory: {TARGET_DIR}\n")
    process_files()
    print("\n✓ Translation complete!")

if __name__ == '__main__':
    main()
