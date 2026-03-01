#!/usr/bin/env python3
"""
Script to translate Spanish files to English and save to EFI_EnglishTraining folder
"""

import os
import shutil
from pathlib import Path
from deep_translator import GoogleTranslator
import json

# Try to import libraries for different file types
try:
    from docx import Document
    from docx.shared import Pt, RGBColor
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# Directory paths
SOURCE_DIR = Path('/workspaces/EFI_Traning')
TARGET_DIR = Path('/workspaces/EFI_EnglishTraining')

# Translation mapping for folder and file names
SPANISH_TO_ENGLISH_NAMES = {
    'Borradores': 'Drafts',
    'Handbook Diario': 'Daily Handbook',
    'Labs': 'Labs',
    'PowerPoint Presentation': 'PowerPoint Presentation',
    'Quizzes': 'Quizzes',
    'Curso 1 borrador remix.docx': 'Course 1 Draft Remix.docx',
    'Curso 1 borrador.docx': 'Course 1 Draft.docx',
    'Dia 1 - EFI basics.docx': 'Day 1 - EFI Basics.docx',
    'Dia 2 -Power_Ground.docx': 'Day 2 - Power_Ground.docx',
    'Dia 3 - Laptops y conexiones.docx': 'Day 3 - Laptops and Connections.docx',
    'Dia 4 - Tiempos de Inyeccion y Dwell.docx': 'Day 4 - Injection Times and Dwell.docx',
    'Dia 5 - Basic Tuning VE_Spark.docx': 'Day 5 - Basic Tuning VE_Spark.docx',
    'Dia 6 - Cranking and AFR Corrections.docx': 'Day 6 - Cranking and AFR Corrections.docx',
    'Dia 7 - Troubleshooting.docx': 'Day 7 - Troubleshooting.docx',
    'Dia 8 - Final Review.docx': 'Day 8 - Final Review.docx',
    'Handbook Dia 1.docx': 'Handbook Day 1.docx',
    'Handbook Dia 2.docx': 'Handbook Day 2.docx',
    'Handbook Dia 3.docx': 'Handbook Day 3.docx',
    'Handbook Dia 4.docx': 'Handbook Day 4.docx',
    'Handbook Dia 5.docx': 'Handbook Day 5.docx',
    'Handbook Dia 6.docx': 'Handbook Day 6.docx',
    'Prompt for Hnadbooks.docx': 'Prompt for Handbooks.docx',
    'Reporte de Laboratorio 1.xlsx': 'Laboratory Report 1.xlsx',
    'Reporte lab 2.xlsx': 'Laboratory Report 2.xlsx',
    'Dia 1.pptx': 'Day 1.pptx',
    'Dia 3 - LA INTERFAZ DIGITAL Y EL DESPERTAR DE LA ECU.pptx': 'Day 3 - The Digital Interface and the Awakening of the ECU.pptx',
    'DÍA 2 — ALIMENTACIÓN, CONTROL DIGITAL Y RUIDO.pptx': 'Day 2 - Power Supply, Digital Control and Noise.pptx',
    'Lectura 1.xlsx': 'Reading 1.xlsx',
    'Quiz Dia 1.docx': 'Quiz Day 1.docx',
    'Quiz Dia 2.docx': 'Quiz Day 2.docx',
    'Prontuario Clase Piloto.docx': 'Pilot Class Handbook.docx',
    'Proposal for Collaboration AlphaX Automeca.docx': 'Proposal for Collaboration AlphaX Automeca.docx',
}

def translate_text(text, max_retries=3):
    """Translate Spanish text to English"""
    if not text or not text.strip():
        return text
    
    try:
        translator = GoogleTranslator(source_language='es', target_language='en')
        # Split text into chunks if too long (limit is 5000 chars per request)
        if len(text) > 4500:
            chunks = []
            words = text.split()
            current_chunk = []
            current_length = 0
            
            for word in words:
                if current_length + len(word) + 1 > 4500:
                    chunks.append(' '.join(current_chunk))
                    current_chunk = [word]
                    current_length = len(word)
                else:
                    current_chunk.append(word)
                    current_length += len(word) + 1
            
            if current_chunk:
                chunks.append(' '.join(current_chunk))
            
            translated = []
            for chunk in chunks:
                try:
                    translated.append(translator.translate(chunk))
                except:
                    translated.append(chunk)  # Fall back to original if translation fails
            return ' '.join(translated)
        else:
            return translator.translate(text)
    except Exception as e:
        print(f"Error translating text: {e}")
        return text

def translate_docx(input_path, output_path):
    """Translate a DOCX file"""
    if not HAS_DOCX:
        print(f"Skipping {input_path} - python-docx not installed")
        return False
    
    try:
        doc = Document(input_path)
        
        # Translate paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                paragraph.text = translate_text(paragraph.text)
        
        # Translate table content
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        if paragraph.text.strip():
                            paragraph.text = translate_text(paragraph.text)
        
        doc.save(output_path)
        print(f"✓ Translated: {input_path}")
        return True
    except Exception as e:
        print(f"✗ Error translating {input_path}: {e}")
        return False

def translate_xlsx(input_path, output_path):
    """Translate an XLSX file"""
    if not HAS_XLSX:
        print(f"Skipping {input_path} - openpyxl not installed")
        return False
    
    try:
        wb = openpyxl.load_workbook(input_path)
        
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value and isinstance(cell.value, str):
                        translated = translate_text(cell.value)
                        cell.value = translated
        
        wb.save(output_path)
        print(f"✓ Translated: {input_path}")
        return True
    except Exception as e:
        print(f"✗ Error translating {input_path}: {e}")
        return False

def translate_pptx(input_path, output_path):
    """Translate a PPTX file"""
    if not HAS_PPTX:
        print(f"Skipping {input_path} - python-pptx not installed")
        return False
    
    try:
        prs = Presentation(input_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip():
                        translated = translate_text(shape.text)
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.text = translated
                
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                cell.text = translate_text(cell.text)
        
        prs.save(output_path)
        print(f"✓ Translated: {input_path}")
        return True
    except Exception as e:
        print(f"✗ Error translating {input_path}: {e}")
        return False

def get_translated_filename(filename):
    """Get the translated version of a filename"""
    return SPANISH_TO_ENGLISH_NAMES.get(filename, filename)

def translate_filename(filename):
    """Attempt to translate filename using the mapping"""
    return get_translated_filename(filename)

def copy_and_translate_files(source_base, target_base, rel_path=''):
    """Recursively copy and translate files"""
    source_path = Path(source_base) / rel_path
    target_path = Path(target_base) / rel_path
    
    if not source_path.exists():
        return
    
    # Create target directory if it doesn't exist
    target_path.mkdir(parents=True, exist_ok=True)
    
    for item in source_path.iterdir():
        item_rel_path = item.relative_to(source_base)
        target_item = target_base / item_rel_path
        
        # Translate the name
        translated_name = translate_filename(item.name)
        target_item = target_item.parent / translated_name
        
        if item.is_dir():
            if item.name not in ['.git', '__pycache__', '.pytest_cache']:
                copy_and_translate_files(source_base, target_base, str(item_rel_path))
        else:
            try:
                if item.suffix.lower() == '.docx':
                    target_item.parent.mkdir(parents=True, exist_ok=True)
                    translate_docx(item, target_item)
                elif item.suffix.lower() == '.xlsx':
                    target_item.parent.mkdir(parents=True, exist_ok=True)
                    translate_xlsx(item, target_item)
                elif item.suffix.lower() == '.pptx':
                    target_item.parent.mkdir(parents=True, exist_ok=True)
                    translate_pptx(item, target_item)
                elif item.suffix.lower() not in ['.pyc', '.git']:
                    # For other files, just copy them
                    target_item.parent.mkdir(parents=True, exist_ok=True)
                    shutil.copy2(item, target_item)
                    print(f"✓ Copied: {item.name}")
            except Exception as e:
                print(f"✗ Error processing {item.name}: {e}")

def main():
    print("Starting translation process...")
    print(f"Source directory: {SOURCE_DIR}")
    print(f"Target directory: {TARGET_DIR}")
    print()
    
    # Create target directory
    TARGET_DIR.mkdir(exist_ok=True)
    
    # Copy and translate files
    copy_and_translate_files(SOURCE_DIR, TARGET_DIR)
    
    print("\nTranslation process completed!")
    print(f"Files saved to: {TARGET_DIR}")

if __name__ == '__main__':
    main()
